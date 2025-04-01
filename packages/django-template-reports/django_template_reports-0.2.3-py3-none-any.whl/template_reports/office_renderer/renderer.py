from pptx import Presentation

from .charts import process_chart
from .paragraphs import process_paragraph
from .tables import process_table_cell


def render_pptx(template, context, output, perm_user):
    """
    Render the PPTX template (a path string or a file-like object) using the provided context and save to output.
    'output' can be a path string or a file-like object. If it's a file-like object, it will be rewound after saving.
    """
    # Support template as a file path or file-like object.
    if isinstance(template, str):
        prs = Presentation(template)
    else:
        template.seek(0)
        prs = Presentation(template)

    errors = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # 1) Process text frames (non-table).
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    # Merge any placeholders that are split across multiple runs.
                    try:
                        process_paragraph(
                            paragraph,
                            context,
                            perm_user=perm_user,
                            mode="normal",  # for text frames
                        )
                    except Exception as e:
                        errors.append(f"Error in paragraph: {e}")
            # 2) Process tables.
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        try:
                            process_table_cell(
                                cell,
                                context,
                                perm_user,
                            )
                        except Exception as e:
                            errors.append(f"Error in table cell: {e}")
            # 3) Process chart spreadsheets.
            if getattr(shape, "has_chart", False):
                try:
                    process_chart(
                        shape.chart,
                        context,
                        perm_user,
                    )
                except Exception as e:
                    errors.append(f"Error in chart: {e}")

    if errors:
        print("Rendering aborted due to the following errors:")
        for err in set(errors):
            print(f" - {err}")
        print("Output file not saved.")
        return None, errors

    # Save to output (file path or file-like object)
    if isinstance(output, str):
        prs.save(output)
    else:
        prs.save(output)
        output.seek(0)

    return output, None
