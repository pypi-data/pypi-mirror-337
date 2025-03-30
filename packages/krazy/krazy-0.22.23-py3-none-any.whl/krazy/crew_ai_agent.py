import chromadb
from chromadb.utils.embedding_functions import EmbeddingFunction
import openai
import os
from crewai import Agent, Task, LLM, Crew, Process
from pathlib import Path
from pydantic import BaseModel
from typing import Optional
import keyring as kr
from tavily import TavilyClient
import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from pydantic import BaseModel, Field, ConfigDict, field_validator
from pypdf import PdfReader
import os
import requests
from bs4 import BeautifulSoup
from crewai.tools import BaseTool
from typing import Type
import pandas as pd
from pathlib import Path
import psycopg2
import keyring as kr
import json
from io import StringIO
from pptx import Presentation  # Install using: pip install python-pptx
from docx import Document  # Install using: pip install python-docx
import fitz  # PyMuPDF
import openai
from datetime import datetime
from typing import List, Dict, Optional
import logging
from tenacity import retry, stop_after_attempt, wait_fixed
from contextlib import contextmanager
import time
import threading
import hashlib
from sentence_transformers import SentenceTransformer, util
from langchain.text_splitter import RecursiveCharacterTextSplitter
import torch


class Config(BaseModel):
    # Azure OpenAI
    azure_open_ai_api_key: str
    azure_version: str
    azure_endpoint: str
    model: str = "gpt-4o"
    model_32k: str = "gpt-4-32k"
    summary_azure_model: str = "gpt-35-turbo-16k"
    embedding_model: str = "text-embedding-ada-002"

    # Tavily
    tavily_api_key: Optional[str] = None

    # PostgreSQL
    postgres_dbname: str = "postgres"
    postgres_user: str = "root"
    postgres_password: Optional[str] = None
    postgres_host: Optional[str] = None
    postgres_port: str = "5432"

    @property
    def endpoint_url(self) -> str:
        return f"{self.azure_endpoint}/openai/deployments/text-embedding-ada-002/completions?api-version=2023-07-01-preview"

    @property
    def embedder_config(self) -> Dict:
        return {
            "embedder": {
                "provider": "azure_openai",
                "config": {
                    "deployment_name": "text-embedding-ada-002",
                    "api_key": self.azure_open_ai_api_key,
                    "api_base": self.endpoint_url,
                },
            }
        }

class AIAssistant:
    def __init__(self,
                 config: Config,
                 db_path="DBs/chroma_memory",
                 embedding_model: str = None,
                 output_schema: Optional[BaseModel] = None,
                 llm_choice: str = 'azure',
                 verbose=True):

        self.config = config
        embedding_model = embedding_model or self.config.embedding_model

        # ✅ Ensure ChromaDB path exists and init memory client
        project_dir = Path(__file__).parent
        absolute_db_path = project_dir.joinpath(db_path)
        os.makedirs(absolute_db_path, exist_ok=True)

        self.memory_client = chromadb.PersistentClient(path=str(absolute_db_path))
        self.memory_collection = self.memory_client.get_or_create_collection(
            name="assistant_memory",
            embedding_function=AzureOpenAIEmbeddingFunction(
                api_key=self.config.azure_open_ai_api_key,
                api_base=self.config.azure_endpoint,
                api_version=self.config.azure_version,
                model=embedding_model
            )
        )

        # output schema
        self.output_schema = output_schema

        # output variables
        self.response_full = None
        self.response = None
        self.response_json = None
        self.response_dict = None
        self.response_pydantic = None

        # ✅ Tools
        self.pdf_reader_tool = PDFReaderTool()
        self.csv_reader_tool = CSVCustomReaderTool()
        self.web_search_tool = WebSearchTool()
        self.postgres_query_tool = PostgresQueryTool()
        self.excel_import_tool = ExcelImportTool()
        self.csv_export_tool = CSVorExcelExportTool()
        self.word_reader_tool = WordFileReaderTool()
        self.powerpoint_reader_tool = PowerPointFileReaderTool()
        self.folder_list_files_tool = FolderListFilesTool()
        self.llm_options = {'azure': LLM(
            model=f"azure/{self.config.model}",
            api_key=self.config.azure_open_ai_api_key,
            api_base=self.config.azure_endpoint,
            api_version=self.config.azure_version,
            temperature=0
        ), 'ollama': LLM(
            model='ollama/deepseek-r1:7b',
            base_url='http://localhost:11434',
            temperature=0
        )}

        # ✅ LLM Setup
        self.llm = self.llm_options[llm_choice]

        # ✅ Agent
        self.agent = Agent(
            role="AI Assistant",
            goal="Answer user queries using data and long-term memory if provided.",
            backstory="You are a knowledgeable assistant with memory capabilities.",
            llm=self.llm,
            verbose=verbose,
        )

        # ✅ Task
        self.task = Task(
            description="{query}",
            expected_output="A response to question based on given information.",
            agent=self.agent,
            tools=[
                self.pdf_reader_tool, self.csv_reader_tool, self.web_search_tool,
                self.postgres_query_tool, self.excel_import_tool, self.csv_export_tool,
                self.word_reader_tool, self.powerpoint_reader_tool, self.folder_list_files_tool
            ],
            output_json=self.output_schema
        )

        # ✅ Crew
        self.crew = Crew(
            name="AI_Memory_Crew",
            agents=[self.agent],
            tasks=[self.task],
            process=Process.sequential
        )

    def prompt_generator(self, query, data, memory):
        if query:
            prompt = str(query)
        else:
            prompt = f'Answer this question:\n Question:{query}.'

        if data:
            prompt += f"\nUse this data to answer the question:\nData: {str(data)}."

        if memory is not None:
            prompt += f"\\Use this context:\nContext: {str(memory)}."

        return prompt

    def store_memory(self, user_input, agent_response):
        self.memory_collection.add(
            ids=[str(len(self.memory_collection.get()["ids"]))],
            documents=[f"User: {user_input}\nAssistant: {agent_response}"]
        )

    def retrieve_memory(self, query, top_k=5):
        results = self.memory_collection.query(query_texts=[query], n_results=top_k)
        return "\n".join(results["documents"][0]) if results["documents"] else "No relevant memories found."

    def invoke(self, prompt, data=None, use_memory=True, save_chat_history=True):
        if use_memory:
            retrieved_memories = self.retrieve_memory(prompt)
        else:
            retrieved_memories = None

        generated_prompt = self.prompt_generator(query=prompt, data=None, memory=retrieved_memories)

        self.response_full = self.crew.kickoff(inputs={"query": generated_prompt})

        if "content_filter" in str(self.response_full):
            return "⚠️ Your request was blocked due to content policy violations. Please rephrase and try again."

        try:
            self.response = self.response_full.raw
        except Exception as e:
            self.response = f"Error: {e}"

        if save_chat_history:
            self.store_memory(prompt, self.response)

        if self.output_schema:
            self.response_pydantic = self.output_schema.model_validate_json(self.response)
            self.response_dict = self.response_pydantic.model_dump()
            self.response_json = self.response_pydantic.model_dump_json()

        return self.response

# ✅ Example Usage
if __name__ == "__main__":
    assistant = AIAssistant()
    query = "Hi, my name is John Doe"
    response = assistant.invoke(query)

    print(f"\n🗨️ User: {query}")
    print(f"🤖 AI Response: {response}")

    print("\n🔍 Checking Memory Storage...")
    print(assistant.retrieve_memory(query))



# ______________________ tools


class WebSearchToolInput(BaseModel):
    """Input model for the WebSearchTool."""
    query: str = Field(..., description="The search query to fetch results from the web.")
    num_results: int = Field(default=3, description="Number of search results to retrieve.")

class WebSearchTool(BaseTool):
    """A CrewAI-compatible tool that searches the web and extracts content from search results."""
    
    name: str = "web_search"
    description: str = "Searches the web using DuckDuckGo and extracts content from results."
    model_config = ConfigDict(arbitrary_types_allowed=True)  # ✅ Allow custom types
    args_schema:Type[BaseModel] = WebSearchToolInput  # ✅ Use the input schema

    def _run(self, query: str, num_results: int = 3) -> str:
        """
        Searches the web and extracts titles, links, and content.

        Args:
            query (str): The search query.
            num_results (int): Number of results to return.

        Returns:
            str: Formatted search results.
        """
        search_url = f"https://duckduckgo.com/html/?q={query.replace(' ', '+')}"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }

        response = requests.get(search_url, headers=headers)
        if response.status_code != 200:
            return f"❌ Error: Unable to fetch search results. Status Code: {response.status_code}"

        soup = BeautifulSoup(response.text, "html.parser")
        extracted_results = []

        # ✅ Extract search results
        for result in soup.find_all("a", class_="result__a", limit=num_results):
            title = result.get_text()
            link = result["href"]
            content = self.fetch_page_content(link)

            extracted_results.append(f"🔹 **{title}**\n🔗 {link}\n📜 {content}\n" + "-" * 50)

        return "\n\n".join(extracted_results) if extracted_results else "❌ No results found."

    def fetch_page_content(self, url: str) -> str:
        """
        Fetches and extracts the main content of a webpage.

        Args:
            url (str): The webpage URL.

        Returns:
            str: Extracted text content.
        """
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }

            response = requests.get(url, headers=headers, timeout=5)
            if response.status_code != 200:
                return f"❌ Error: Unable to fetch page. Status Code: {response.status_code}"

            soup = BeautifulSoup(response.text, "html.parser")
            paragraphs = soup.find_all("p")
            page_content = "\n".join([p.get_text() for p in paragraphs])

            return page_content[:1000] + "..." if len(page_content) > 1000 else page_content  # Limit to 1000 chars

        except Exception as e:
            return f"❌ Error fetching content: {e}"

class WebSearchToolTavilyInput[BaseModel]:
    query: str = Field(..., description="The search query to fetch results from the web.")

class WebSearchToolTavily(BaseTool):  # Ensure it inherits from BaseTool
    name: str = "web_search"
    description: str = "Searches the web and returns results."
    model_config = ConfigDict(arbitrary_types_allowed=True)  # Use ConfigDict for Pydantic v2
    args_schema:Type[BaseModel]=WebSearchToolTavilyInput

    def _run(self, query: str) -> str:  # Implement the required _run method
        """
        Searches the web for the given query.
        Args:
            query (str): The search query.
        Returns:
            str: The search results as a formatted string.
        """
        try:
            config = Config()
            web_search_tool = TavilyClient(api_key=config.tavily_api_key)
            response = web_search_tool.search(f'{query}', max_results=2)
            print(f'Tavily Response: {response['results']}')
            # Convert response from list to a formatted string
            if response:
                formatted_response = "\n".join([f"- {result}" for result in response['results']])
                return f"Here are the top results for your query:\n{formatted_response}"
            else:
                return "No results found for your query. Please try a different query."
        except Exception as e:
            return f"An error occurred during the web search: {e}"
        

class PDFReaderToolInput(BaseModel):
    """Input schema for the PDF Reader Tool."""
    file_path: str = Field(..., description="Path to the PDF file.")

class PDFReaderTool(BaseTool):
    name: str = "pdf_reader"
    description: str = "Reads and extracts text from a PDF file."
    model_config = ConfigDict(arbitrary_types_allowed=True)  # Use ConfigDict for Pydantic v2
    args_schema:Type[BaseModel]=PDFReaderToolInput

    def _run(self, file_path:str) -> str:

        """
        Extracts text from all pages of a PDF.

        Args:
            file_path (str): Path to the PDF file.

        Returns:
            str: Full text content of the PDF.
        """
        if not Path(file_path).exists():
            raise FileNotFoundError(f"❌ File does not exist: {file_path}")

        try:
            text = ""
            with fitz.open(file_path) as pdf:
                for page_num in range(len(pdf)):
                    page = pdf[page_num]
                    text += page.get_text()
            return text.strip()
        
        except Exception as e:
            return f"❌ Error reading PDF: {e}"

class CSVCustomReaderToolInput(BaseModel):
    """Input schema for the CSV Reader Tool."""
    code: str = Field(..., description="Python code for extracting information from the CSV.")
    file_path: str = Field(..., description="Path to the CSV file.")

class CSVCustomReaderTool(BaseTool):
    """Tool for reading and executing agent-generated Python code on a CSV file."""
    
    name: str = "csv_reader"
    description: str = "Executes Python code on a CSV file and returns processed results. Include result = ... in the code. result is what will be returned."
    args_schema:Type[BaseModel] = CSVCustomReaderToolInput

    def _run(self, code: str, file_path: str=None) -> str:
        """Executes the agent-provided Python code on the given CSV file."""
        try:
            if file_path is None:
                return "Provide path to CSV file"
            if Path(file_path).exists:
                pass
            else:
                return 'File does not exist'
            
            df = pd.DataFrame()
            exec_globals = {"df": df}
            exec(code, exec_globals)
            
            # The agent should always return the result in a variable named `result`
            result = exec_globals.get("result", "❌ No valid result found.")

            return result
        
        except Exception as e:
            return f"❌ Error executing code: {e}"


class PostgresQueryInput(BaseModel):
    """Schema for Postgres Query Tool."""
    query: str = Field(..., description="Select SQL query to execute in PostgreSQL.")

class PostgresQueryTool(BaseTool):
    """Tool for querying a PostgreSQL database."""
    
    name: str = "postgres_query"
    description: str = "Executes Select SQL queries on a PostgreSQL database."
    args_schema:Type[BaseModel] = PostgresQueryInput

    def _run(self, query: str):
        """Executes a SQL query and returns the result as JSON."""
        try:
            if query.lower().startswith("select"):
                pass
            else:
                # print(f'Query: {query}')
                return "Only SELECT queries are allowed."
            
            conn = psycopg2.connect(
                dbname='postgres',
                user="root",
                password=kr.get_password('aws-posgtres', 'root'),
                host=kr.get_password('aws-posgtres', 'endpoint'),
                port="5432"
            )
            cursor = conn.cursor()
            cursor.execute(query)
            records = cursor.fetchall()
            col_names = [desc[0] for desc in cursor.description]  # Extract column names
            cursor.close()
            conn.close()
            return json.dumps([dict(zip(col_names, row)) for row in records])  # Convert to JSON
        except Exception as e:
            return f"❌ Error executing query: {e}"

class ExcelImportToolInput(BaseModel):
    """Input schema for importing and processing an Excel file."""
    file_path: str = Field(..., description="Path to the Excel file to import.")
    code: str = Field(..., description="Python code to execute on the DataFrame.")

class ExcelImportTool(BaseTool):
    """Tool for importing an Excel file, executing LLM-generated code on it, and returning results."""
    
    name: str = "excel_import"
    description: str = "Imports an Excel file, processes it using the provided Python code, and returns the result."
    args_schema: Type[BaseModel] = ExcelImportToolInput

    def _run(self, file_path: str, code: str):
        """Loads an Excel file into a DataFrame, executes the given Python code, and returns the output."""
        try:
            if not os.path.exists(file_path):
                return f"❌ Error: The file '{file_path}' does not exist."

            xl = pd.ExcelFile(file_path)  # Loads the excel file
            exec_globals = {"xl": xl}  # Define execution environment

            exec(code, exec_globals)  # Execute LLM-provided code

            # Extract result from the execution
            result = exec_globals.get("result", "❌ No valid result found. Ensure LLM sets 'result' variable.")

            return result  # Return processed result

        except Exception as e:
            return f"❌ Error processing Excel file: {e}"

class CSVorExcelExportToolInput(BaseModel):
    """Input schema for exporting data to CSV."""
    file_path: str = Field(..., description="Path to save the exported CSV file.")
    dataframe_json: str = Field(..., description="DataFrame in JSON format to export.")
    export_as: str = Field("csv", description="Export format (default: csv). Mention xlsx for excel export.")

class CSVorExcelExportTool(BaseTool):
    """Tool for exporting a pandas DataFrame to a CSV file."""
    
    name: str = "csv_export"
    description: str = "Exports data (from a DataFrame) to a CSV file."
    args_schema: Type[BaseModel] = CSVorExcelExportToolInput

    def _run(self, file_path: str, dataframe_json: str, export_as: str = "csv"):
        """Exports the provided JSON data to a CSV file."""
        try:
            if export_as.lower() == "csv":
                df = pd.read_json(StringIO(dataframe_json))   # Convert JSON to DataFrame
                df.to_csv(file_path, index=False)  # Save as CSV
            if export_as.lower() == "xlsx":
                df = pd.read_json(StringIO(dataframe_json))
                df.to_excel(file_path, index=False)

            return f"✅ Data exported successfully to {file_path}"

        except Exception as e:
            return f"❌ Error exporting CSV: {e}"


class WordFileReaderToolInput(BaseModel):
    """Input schema for reading a Word (.docx) file."""
    file_path: str = Field(..., description="Path to the Word file.")

class WordFileReaderTool(BaseTool):
    """Tool to read and extract text from a Microsoft Word (.docx) file."""
    
    name: str = "word_reader"
    description: str = "Reads and extracts text from a Word document (.docx)."
    args_schema: Type[BaseModel] = WordFileReaderToolInput

    def _run(self, file_path: str):
        """Reads and extracts text from a Word file."""
        try:
            if not os.path.exists(file_path):
                return f"❌ Error: The file '{file_path}' does not exist."

            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

            return text.strip() if text else "❌ No text found in the document."

        except Exception as e:
            return f"❌ Error reading Word file: {e}"


class PowerPointFileReaderToolInput(BaseModel):
    """Input schema for reading a PowerPoint (.pptx) file."""
    file_path: str = Field(..., description="Path to the PowerPoint file.")

class PowerPointFileReaderTool(BaseTool):
    """Tool to read and extract text from a PowerPoint (.pptx) file."""
    
    name: str = "pptx_reader"
    description: str = "Reads and extracts text from a PowerPoint presentation (.pptx)."
    args_schema: Type[BaseModel] = PowerPointFileReaderToolInput

    def _run(self, file_path: str):
        """Reads and extracts text from a PowerPoint file."""
        try:
            if not os.path.exists(file_path):
                return f"❌ Error: The file '{file_path}' does not exist."

            ppt = Presentation(file_path)
            slides_text = []

            for slide in ppt.slides:
                slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                slides_text.append(slide_text)

            full_text = "\n\n".join(slides_text)

            return full_text.strip() if full_text else "❌ No text found in the presentation."

        except Exception as e:
            return f"❌ Error reading PowerPoint file: {e}"

class FolderListFilesToolInput(BaseModel):
    """Input schema for listing files in a folder."""
    folder_path: str = Field(..., description="Path to the folder whose files need to be listed.")

class FolderListFilesTool(BaseTool):
    """Tool to list all files in a specified folder."""
    
    name: str = "folder_list_files"
    description: str = "Lists all files in a given folder."
    args_schema: Type[BaseModel] = FolderListFilesToolInput

    def _run(self, folder_path: str):
        """Lists all files in a folder."""
        try:
            if not os.path.exists(folder_path):
                return f"❌ Error: The folder '{folder_path}' does not exist."

            if not os.path.isdir(folder_path):
                return f"❌ Error: The path '{folder_path}' is not a directory."

            files = os.listdir(folder_path)
            file_list = [f for f in files if os.path.isfile(os.path.join(folder_path, f))]

            return file_list if file_list else "📁 No files found in the folder."

        except Exception as e:
            return f"❌ Error listing files: {e}"


class RAGPostgres:
    """
    A Retrieval-Augmented Generation (RAG) engine that uses PostgreSQL for vector and document storage.

    This class facilitates:
    - Connecting to an AWS-hosted PostgreSQL database.
    - Ingesting and chunking text data from PDF, Word, PowerPoint, Markdown, and Text files.
    - Computing embeddings using Azure OpenAI.
    - Storing embeddings and metadata into a PostgreSQL table with full-text search capability.
    - Detecting and preventing re-ingestion of the same file using checksum.
    - Watching a folder for new files and ingesting them automatically.
    - Searching similar documents using vector similarity and reranking them with sentence-level cosine similarity.
    - Highlighting the most relevant sentence from each chunk during search.

    Args:
        doc_table (str): Name of the table where documents are stored.
        conn_param (dict): Dictionary containing PostgreSQL connection parameters such as 'dbname', 'user', 'password', 'host'.

    Attributes:
        schema (str): Schema in which all RAG-related tables are stored.
        doc_table (str): Table for document chunks with embeddings.
        chat_table (str): Table for storing chat history (not used in current scope).
        ingested_files_table (str): Table to track already ingested files via checksum.
        watcher_thread (threading.Thread): Thread object for watching file system.
        stop_watcher_flag (threading.Event): Event used to gracefully stop the watcher thread.
        client (openai.AzureOpenAI): Azure OpenAI client for embeddings.
        sentence_encoder (SentenceTransformer): Local sentence transformer for reranking.
    """
    def __init__(self, doc_table='rag_documents', conn_param:dict=None):
        self.config = Config()
        self.conn_kwargs = conn_param
        if self.conn_kwargs:
            self.conn = self.connect_to_postgres_aws()
        else:
            raise ValueError("Connection kwargs required.")
        
        self.schema = 'ai'
        self.doc_table = doc_table
        self.chat_table = 'chat_history'
        self.ingested_files_table = 'ingested_files'
        self.watcher_thread = None
        self.stop_watcher_flag = threading.Event()

        self.client = openai.AzureOpenAI(
            api_key=self.config.azure_open_ai_api_key,
            azure_endpoint=self.config.azure_endpoint,
            api_version="2023-07-01-preview"
        )
        self.sentence_encoder = SentenceTransformer("all-mpnet-base-v2")

    def connect_to_postgres_aws(self):
        return psycopg2.connect(
            dbname=self.conn_kwargs.get('dbname', 'postgres'),
            user=self.conn_kwargs.get('user', 'root'),
            password=self.conn_kwargs.get('password'),
            host=self.conn_kwargs.get('host'),
            port=self.conn_kwargs.get('port', '5432')
        )

    @contextmanager
    def get_cursor(self):
        with self.conn.cursor() as cur:
            yield cur

    def setup_tables(self):
        with self.get_cursor() as cur:
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {self.schema}.{self.doc_table} (
                    id SERIAL PRIMARY KEY,
                    content TEXT,
                    metadata JSONB,
                    embedding VECTOR(1536),
                    tsv TSVECTOR
                );
            """)
            cur.execute(f"""
                CREATE INDEX IF NOT EXISTS idx_embedding_vector ON {self.schema}.{self.doc_table} USING ivfflat (embedding vector_cosine_ops);
            """)
            cur.execute(f"""
                CREATE INDEX IF NOT EXISTS idx_metadata ON {self.schema}.{self.doc_table} USING GIN (metadata);
            """)
            cur.execute(f"""
                CREATE INDEX IF NOT EXISTS idx_tsv ON {self.schema}.{self.doc_table} USING GIN(tsv);
            """)
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {self.schema}.{self.chat_table} (
                    id SERIAL PRIMARY KEY,
                    chat_agent TEXT,
                    agent_name TEXT,
                    user_input TEXT,
                    assistant_response TEXT,
                    timestamp TIMESTAMPTZ DEFAULT NOW(),
                    session_id UUID DEFAULT gen_random_uuid(),
                    embedding VECTOR(1536)
                );
            """)
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {self.schema}.{self.ingested_files_table} (
                    id SERIAL PRIMARY KEY,
                    file_name TEXT UNIQUE,
                    checksum TEXT
                );
            """)
        self.conn.commit()

    def calculate_checksum(self, file_path: str) -> str:
        hasher = hashlib.sha256()
        with open(file_path, 'rb') as f:
            while chunk := f.read(8192):
                hasher.update(chunk)
        return hasher.hexdigest()

    def is_file_already_ingested(self, file_path: str) -> bool:
        checksum = self.calculate_checksum(file_path)
        with self.get_cursor() as cur:
            cur.execute(
                f"SELECT 1 FROM {self.schema}.{self.ingested_files_table} WHERE file_name = %s AND checksum = %s",
                (Path(file_path).name, checksum)
            )
            return cur.fetchone() is not None

    def record_file_ingestion(self, file_path: str):
        checksum = self.calculate_checksum(file_path)
        with self.get_cursor() as cur:
            cur.execute(
                f"INSERT INTO {self.schema}.{self.ingested_files_table} (file_name, checksum) VALUES (%s, %s) ON CONFLICT DO NOTHING",
                (Path(file_path).name, checksum)
            )
        self.conn.commit()

    @retry(stop=stop_after_attempt(3), wait=wait_fixed(2))
    def get_embedding(self, text: str) -> List[float]:
        response = self.client.embeddings.create(
            input=[text],
            model="text-embedding-ada-002"
        )
        return response.data[0].embedding

    def chunk_text(self, text: str, max_chunk_tokens=400, overlap=40) -> List[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=max_chunk_tokens,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", ".", " "],
            length_function=lambda x: len(x.split())
        )
        return splitter.split_text(text)

    def extract_text_from_pdf(self, file_path: str) -> str:
        with fitz.open(file_path) as pdf:
            return "\n".join(page.get_text() for page in pdf).strip()

    def extract_text_from_word(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

    def extract_text_from_ppt(self, file_path: str) -> str:
        prs = Presentation(file_path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

    def extract_text_from_txt_or_md(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read().strip()

    def ingest_file(self, file_path: str, metadata: Optional[Dict] = None):
        if self.is_file_already_ingested(file_path):
            print(f"🔁 Skipping already ingested file: {file_path}")
            return

        ext = Path(file_path).suffix.lower()
        metadata = metadata or {"source": str(file_path)}
        try:
            if ext == ".pdf":
                text = self.extract_text_from_pdf(file_path)
            elif ext == ".docx":
                text = self.extract_text_from_word(file_path)
            elif ext == ".pptx":
                text = self.extract_text_from_ppt(file_path)
            elif ext in [".txt", ".md"]:
                text = self.extract_text_from_txt_or_md(file_path)
            else:
                print(f"❌ Unsupported file type for ingestion: {ext}")
                return

            chunks = self.chunk_text(text)
            self.insert_document_chunks(chunks, metadata)
            self.record_file_ingestion(file_path)
        except Exception as e:
            print(f"❌ Failed to ingest file {file_path}: {e}")

    def insert_document_chunks(self, chunks: List[str], metadata: Dict):
        with self.get_cursor() as cur:
            for chunk in chunks:
                embedding = self.get_embedding(chunk)
                cur.execute(
                    f"""
                    INSERT INTO {self.schema}.{self.doc_table} (content, metadata, embedding, tsv)
                    VALUES (%s, %s, %s, to_tsvector(%s))
                    """,
                    (chunk, json.dumps(metadata), embedding, chunk)
                )
        self.conn.commit()

    def watch_folder(self, folder_path: str, polling_interval: int = 10):
        print(f"👀 Watching folder: {folder_path}")
        folder = Path(folder_path)

        def watch():
            while not self.stop_watcher_flag.is_set():
                for file_path in folder.glob("*"):
                    if file_path.is_file():
                        self.ingest_file(str(file_path))
                time.sleep(polling_interval)

        self.watcher_thread = threading.Thread(target=watch, daemon=True)
        self.watcher_thread.start()

    def stop_watcher(self):
        if self.watcher_thread:
            self.stop_watcher_flag.set()
            self.watcher_thread.join()
            print("🛑 Watcher stopped.")

    def highlight_relevant_passage(self, query: str, chunk: str) -> str:
        sentences = [s.strip() for s in chunk.split('.') if len(s.strip()) > 0]
        if not sentences:
            return chunk
        query_embedding = self.sentence_encoder.encode(query, convert_to_tensor=True)
        sentence_embeddings = self.sentence_encoder.encode(sentences, convert_to_tensor=True)
        scores = util.pytorch_cos_sim(query_embedding, sentence_embeddings)[0]
        top_idx = scores.argmax().item()
        return sentences[top_idx]

    def search_similar_documents(self, query: str, top_k: int = 5, metadata_filter: Optional[Dict] = None, highlight=True):
        embedding = self.get_embedding(query)
        sql = f"SELECT content, metadata FROM {self.schema}.{self.doc_table}"
        params = []

        where_clauses = []
        if metadata_filter:
            where_clauses.append("metadata @> %s::jsonb")
            params.append(json.dumps(metadata_filter))
        where_clause = " WHERE " + " AND ".join(where_clauses) if where_clauses else ""

        sql += where_clause
        sql += f" ORDER BY embedding <-> %s::vector LIMIT %s"
        params.extend([embedding, top_k * 3])

        with self.get_cursor() as cur:
            cur.execute(sql, params)
            results = cur.fetchall()

        ranked = self.rerank_results(query, results, final_k=top_k)
        return ranked if highlight else results

    def rerank_results(self, query: str, chunks: List[tuple], final_k: int = 5):
        contents = [c[0] for c in chunks]
        meta = [c[1] for c in chunks]
        query_embedding = self.sentence_encoder.encode(query, convert_to_tensor=True)
        doc_embeddings = self.sentence_encoder.encode(contents, convert_to_tensor=True)
        scores = util.pytorch_cos_sim(query_embedding, doc_embeddings)[0]
        top_results = torch.topk(scores, final_k)
        return [
            {"highlight": contents[i], "metadata": meta[i]}
            for i in top_results.indices.tolist()
        ]
